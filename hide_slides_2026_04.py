from pptx import Presentation
import sys, os

KEEP = [
    "Azure SRE Agent","Priority Processing in Microsoft Foundry",
    "Foundry Agent Service","GPT-5.4",
    "Fabric Mirroring","MySQL","Item Recovery in Microsoft Fabric",
    "Cilium mTLS","Container network logs","Container network metrics",
    "AI Agent for container networking","AKS managed GPU metrics",
    "meshless Istio","Blue-green agent pool","Cross-cluster networking",
    "Kubernetes Application Network",
    "PostgreSQL dashboards with Grafana","EDB workloads","Cosmos DB Mirroring",
    "Transparent Data Encryption","Azure SQL",
    "Azure Developer CLI","azd",
    "Evaluations, Monitoring, and Tracing","OpenTelemetry",
    "Azure Storage Mover","AWS S3",
    "Draft & Deploy on Azure Firewall",
    "Service Retirement","Retirement Book",
    "Technical Update Briefing","Unlocking","Updates by Azure",
    "Agenda","Welcome","Thank",
]

def should_show(slide):
    title = ""
    if slide.shapes.title:
        title = slide.shapes.title.text
    if not title:
        for s in slide.shapes:
            if s.has_text_frame:
                title = s.text_frame.text; break
    tl = title.lower()
    return any(k.lower() in tl for k in KEEP)

src = sys.argv[1]
dst = os.path.splitext(src)[0] + "_filtered.pptx"
prs = Presentation(src)
shown, hidden = 0, 0
for i, slide in enumerate(prs.slides, 1):
    show = should_show(slide)
    el = slide._element
    if show:
        el.attrib.pop("show", None)
        shown += 1
    else:
        el.attrib["show"] = "0"
        hidden += 1
prs.save(dst)
print(f"완료: 표시 {shown}장, 숨김 {hidden}장 → {dst}")