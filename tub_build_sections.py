"""
2026-04 Azure TUB — Voting 선정 슬라이드만 남기고 Area별 Section 구성
=====================================================================
동작:
  1. 음영처리된 25개 항목 외 슬라이드 삭제 (숨김 아님)
  2. 남은 슬라이드를 Area 순서로 재정렬: IH → CH → YJ → 나머지
  3. 각 Area를 PowerPoint Section으로 표시

사용법:
  pip install python-pptx lxml
  python tub_build_sections.py "2026_04 - Azure-TUB - 복사본.pptx"

출력:
  2026_04 - Azure-TUB - 복사본_sectioned.pptx  (원본과 같은 폴더)
"""

import sys, os, uuid
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

# ─────────────────────────────────────────────────────────────
# Voting Dashboard 음영 항목 → (매칭 키워드, Area)
# ─────────────────────────────────────────────────────────────
ITEMS = [
    # (keywords_any_of,  area)
    (["SRE Agent"],                                              "IH"),
    (["Priority Processing", "Performance-Sensitive AI"],        "IH"),
    (["Foundry Agent Service"],                                  "IH"),
    (["GPT-5.4", "GPT 5.4"],                                    "IH"),

    (["Fabric Mirroring", "MySQL", "FabCon"],                   "나머지"),
    (["Item Recovery in Microsoft Fabric"],                      "나머지"),

    (["Cilium mTLS"],                                           "CH"),
    (["Container network logs"],                                 "CH"),
    (["Container network metrics"],                              "CH"),
    (["AI Agent for container networking"],                      "CH"),
    (["AKS managed GPU metrics"],                               "CH"),
    (["meshless Istio"],                                        "CH"),
    (["Blue-green agent pool"],                                 "YJ"),
    (["Cross-cluster networking", "Kubernetes Fleet Manager"],   "CH"),
    (["Kubernetes Application Network"],                        "CH"),

    (["PostgreSQL dashboards", "Grafana"],                      "나머지"),
    (["EDB workloads", "PostgreSQL migration service"],         "나머지"),
    (["Cosmos DB Mirroring"],                                   "나머지"),
    (["Transparent Data Encryption", "Versionless keys"],       "나머지"),
    (["Azure Developer CLI", "azd"],                            "나머지"),

    (["Evaluations, Monitoring, and Tracing"],                  "YJ"),
    (["Ingest OTLP", "OpenTelemetry Collector"],                "YJ"),
    (["Azure Storage Mover", "AWS S3"],                        "YJ"),
    (["Draft & Deploy", "Azure Firewall"],                      "YJ"),
    (["Service Retirement Book", "Retirement Book"],            "YJ"),
]

# 항상 유지할 슬라이드 키워드 (표지, 목차 등)  → area = "커버"
COVER_KEYWORDS = [
    "Technical Update Briefing", "Azure TUB", "Unlocking Innovations",
    "Updates by Azure Product", "Agenda", "Welcome", "Thank",
]

SECTION_ORDER = ["IH", "CH", "YJ", "나머지"]


# ─────────────────────────────────────────────────────────────
# 헬퍼
# ─────────────────────────────────────────────────────────────
def get_title(slide):
    if slide.shapes.title:
        t = slide.shapes.title.text.strip()
        if t:
            return t
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                return t[:300]
    return ""


def classify(title):
    """Return area string, or None if slide should be deleted."""
    tl = title.lower()
    # 커버 슬라이드 유지
    for kw in COVER_KEYWORDS:
        if kw.lower() in tl:
            return "커버"
    # 항목 매칭
    for keywords, area in ITEMS:
        if any(k.lower() in tl for k in keywords):
            return area
    return None  # 삭제 대상


# ─────────────────────────────────────────────────────────────
# 슬라이드 삭제 (python-pptx 공식 미지원 → XML 직접 조작)
# ─────────────────────────────────────────────────────────────
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def delete_slides(prs, indices_to_delete):
    """indices_to_delete: set of 0-based indices to remove."""
    sldIdLst = prs.slides._sldIdLst
    # 역순 삭제 (인덱스 안 밀리도록)
    all_sldIds = list(sldIdLst)
    for idx in sorted(indices_to_delete, reverse=True):
        el = all_sldIds[idx]
        rId = el.get(f"{{{REL_NS}}}id")
        sldIdLst.remove(el)
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


# ─────────────────────────────────────────────────────────────
# 슬라이드 재정렬
# ─────────────────────────────────────────────────────────────
def reorder_slides(prs, new_order):
    """new_order: list of current 0-based indices in desired order."""
    sldIdLst = prs.slides._sldIdLst
    saved = list(sldIdLst)
    for el in saved:
        sldIdLst.remove(el)
    for idx in new_order:
        sldIdLst.append(saved[idx])


# ─────────────────────────────────────────────────────────────
# Section 추가
# ─────────────────────────────────────────────────────────────
P14_NS  = "http://schemas.microsoft.com/office/powerpoint/2010/main"
EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"


def add_sections(prs, section_slide_map):
    """
    section_slide_map: OrderedDict/list of (section_name, [slide_indices])
    slide_indices are 0-based indices AFTER reordering.
    """
    prs_elem = prs.slides._sldIdLst.getparent()  # <p:presentation>

    # 기존 sectionLst 제거
    for ext in prs_elem.findall(f".//{{{P14_NS}}}sectionLst"):
        ext.getparent().getparent().remove(ext.getparent())

    # extLst 확보
    extLst_tag = qn("p:extLst")
    extLst = prs_elem.find(extLst_tag)
    if extLst is None:
        extLst = etree.SubElement(prs_elem, extLst_tag)

    # <p:ext uri="...">
    ext_el = etree.SubElement(extLst, qn("p:ext"))
    ext_el.set("uri", EXT_URI)

    # <p14:sectionLst>
    sectionLst = etree.SubElement(
        ext_el,
        f"{{{P14_NS}}}sectionLst",
        nsmap={"p14": P14_NS},
    )

    # 현재 슬라이드 ID 목록 (정렬된 상태)
    slide_ids = [
        prs.slides._sldIdLst[i].get("id")
        for i in range(len(prs.slides))
    ]

    for sec_name, indices in section_slide_map:
        if not indices:
            continue
        sec_el = etree.SubElement(
            sectionLst,
            f"{{{P14_NS}}}section",
        )
        sec_el.set("name", sec_name)
        sec_el.set("id", "{" + str(uuid.uuid4()).upper() + "}")

        inner_lst = etree.SubElement(sec_el, f"{{{P14_NS}}}sldIdLst")
        for idx in indices:
            sldId_el = etree.SubElement(inner_lst, f"{{{P14_NS}}}sldId")
            sldId_el.set("id", slide_ids[idx])


# ─────────────────────────────────────────────────────────────
# 메인 처리
# ─────────────────────────────────────────────────────────────
def process(src, dst=None):
    if dst is None:
        base, ext = os.path.splitext(src)
        dst = base + "_sectioned" + ext

    print(f"\n파일 로딩: {src}")
    prs = Presentation(src)

    total = len(prs.slides)
    print(f"전체 슬라이드: {total}장\n")

    # 1. 분류
    keep_map = {}   # {original_index: area}
    delete_idx = set()

    print(f"{'슬라이드':>5}  {'분류':^8}  제목")
    print("-" * 90)
    for i, slide in enumerate(prs.slides):
        title = get_title(slide)
        area  = classify(title)
        if area is None:
            delete_idx.add(i)
            print(f"  {i+1:>3}   {'[삭제]':^8}  {title[:70]}")
        else:
            keep_map[i] = area
            print(f"  {i+1:>3}   {('[' + area + ']'):^8}  {title[:70]}")

    print("-" * 90)
    print(f"유지: {len(keep_map)}장  /  삭제: {len(delete_idx)}장\n")

    # 2. 삭제
    delete_slides(prs, delete_idx)

    # 3. 재정렬: 커버 → IH → CH → YJ → 나머지
    # 삭제 후 새로운 인덱스 재계산
    kept_original = sorted(keep_map.keys())           # 원본 인덱스 순서
    new_index_map = {orig: new for new, orig in enumerate(kept_original)}

    area_buckets = {"커버": [], "IH": [], "CH": [], "YJ": [], "나머지": []}
    for orig_idx in kept_original:
        area = keep_map[orig_idx]
        new_i = new_index_map[orig_idx]
        area_buckets[area].append(new_i)

    # 원하는 섹션 순서대로 재정렬
    ordered_keys = ["커버"] + SECTION_ORDER
    new_order = []
    for key in ordered_keys:
        new_order.extend(area_buckets[key])

    reorder_slides(prs, new_order)

    # 4. 재정렬 후 각 섹션의 슬라이드 위치 재계산
    position = 0
    section_slide_map = []
    for key in SECTION_ORDER:   # 커버 섹션은 제외하고 IH/CH/YJ/나머지만 섹션으로
        count = len(area_buckets[key])
        if count:
            cover_count = len(area_buckets["커버"])
            # 커버는 섹션에 포함 안 함 → offset 적용
            start = cover_count + sum(
                len(area_buckets[k]) for k in SECTION_ORDER[:SECTION_ORDER.index(key)]
            )
            indices = list(range(start, start + count))
            section_slide_map.append((key, indices))

    # 5. Section XML 삽입
    add_sections(prs, section_slide_map)

    # 6. 저장
    prs.save(dst)

    print("\n=== 결과 요약 ===")
    for key in ordered_keys:
        cnt = len(area_buckets[key])
        if cnt:
            label = "커버 (섹션 미포함)" if key == "커버" else f"섹션: {key}"
            print(f"  {label}: {cnt}장")
    print(f"\n저장 완료: {dst}\n")
    print("PowerPoint에서 열면 왼쪽 패널에 IH / CH / YJ / 나머지 섹션이 표시됩니다.")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("사용법: python tub_build_sections.py <파일.pptx>")
        sys.exit(1)

    src = sys.argv[1]
    dst = sys.argv[2] if len(sys.argv) > 2 else None

    if not os.path.exists(src):
        print(f"파일 없음: {src}")
        sys.exit(1)

    process(src, dst)