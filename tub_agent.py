"""
TUB Section Agent
=================
Voting Dashboard Excel에서 선정 항목을 읽어 PPTX를 자동 처리합니다.

처리 내용:
  1. 비선정 슬라이드 삭제 (숨김 아님)
  2. Section 구성:  Opening → IH → CH → YJ → MH → Closeout
  3. 카테고리 헤더 슬라이드는 해당 카테고리에 선정 항목이 있으면 보존
  4. Agenda까지는 Opening, Thank You 슬라이드는 Closeout으로 이동

사용법:
  pip install python-pptx openpyxl lxml

  # Excel 자동 감지 (최근 _result 시트) + PPTX 지정
  python tub_agent.py "2026_04 - Azure-TUB - 복사본.pptx"

  # Excel 시트 직접 지정
  python tub_agent.py "2026_04 - Azure-TUB - 복사본.pptx" \
        --excel "Azure_TUB_Voting Dashboard.xlsx" \
        --sheet "2026-04_result"

출력:
  원본파일명_sectioned.pptx (같은 폴더)
"""

import sys
import os
import uuid
import argparse
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
import openpyxl


# ═══════════════════════════════════════════════════════════════
# 1. TUB 표준 카테고리 목록
#    (PPTX의 카테고리 구분 슬라이드 제목과 대조)
# ═══════════════════════════════════════════════════════════════
KNOWN_CATEGORIES = {
    'AI + Machine Learning', 'Analytics', 'Compute', 'Containers',
    'Databases', 'DevOps', 'Developer Tools', 'Hybrid + Multicloud',
    'Identity', 'Integration', 'IoT', 'Management and Governance',
    'Migration', 'Networking', 'Security', 'Storage', 'Web',
    'Retirement', 'Azure Retirement',
}

# Section 표시 순서
SECTION_ORDER = ['IH', 'CH', 'YJ', 'MH']

# 네임스페이스
REL_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
P14_NS = 'http://schemas.microsoft.com/office/powerpoint/2010/main'
EXT_URI = '{521415D9-36F7-43E2-AB2F-B90AF26B5E84}'


# ═══════════════════════════════════════════════════════════════
# 2. Excel 파싱
# ═══════════════════════════════════════════════════════════════

def _is_shaded(cell):
    fill = cell.fill
    if not fill or not fill.fill_type or fill.fill_type == 'none':
        return False
    fg = fill.fgColor
    try:
        rgb = fg.rgb
        return rgb not in {'00000000', 'FFFFFFFF', '00FFFFFF', 'FF000000'}
    except Exception:
        try:
            _ = fg.theme
            return True
        except Exception:
            return fill.fill_type == 'solid'


def load_voting_data(excel_path, sheet_name=None):
    """
    Returns:
      selected   : {normalized_title: {'title', 'area', 'category'}}
      cat_areas  : {category: dominant_area}   (카테고리의 대표 Area)
    """
    wb = openpyxl.load_workbook(excel_path)

    if sheet_name is None:
        candidates = sorted(
            s for s in wb.sheetnames
            if '_result' in s and 'source' not in s
        )
        if not candidates:
            raise ValueError("_result 시트를 찾을 수 없습니다.")
        sheet_name = candidates[-1]
        print(f"[Excel] 시트: {sheet_name}")

    ws = wb[sheet_name]

    selected = {}       # normalized_title → meta
    area_counts = {}    # {category: {area: count}}
    current_cat = ''

    for row in ws.iter_rows(min_row=4, max_row=ws.max_row):
        cell_a = row[0]
        cell_c = row[2] if len(row) > 2 else None

        if not cell_a.value:
            continue

        title = str(cell_a.value).strip().rstrip()
        # 제목 끝의 trailing whitespace/nbsp 제거
        title = ' '.join(title.split())

        # 카테고리 헤더 감지
        if title in KNOWN_CATEGORIES:
            current_cat = title
            area_counts.setdefault(current_cat, {})
            continue

        # 선정 항목 (음영 있음)
        if _is_shaded(cell_a):
            area = ''
            if cell_c and cell_c.value and not isinstance(cell_c.value, (int, float)):
                area = str(cell_c.value).strip()
            if not area:
                area = 'MH'

            norm = _normalize(title)
            selected[norm] = {
                'title': title,
                'area':  area,
                'category': current_cat,
            }
            # 카테고리 area 빈도 집계
            area_counts.setdefault(current_cat, {})
            area_counts[current_cat][area] = area_counts[current_cat].get(area, 0) + 1

    # 카테고리별 대표 Area 결정 (최다 항목 기준, 동점이면 SECTION_ORDER 우선)
    cat_areas = {}
    for cat, counts in area_counts.items():
        if not counts:
            continue
        dominant = max(counts, key=lambda a: (counts[a], -SECTION_ORDER.index(a) if a in SECTION_ORDER else -99))
        cat_areas[cat] = dominant

    return selected, cat_areas


def _normalize(title):
    """공백·대소문자 정규화"""
    return ' '.join(title.lower().split())


# ═══════════════════════════════════════════════════════════════
# 3. 슬라이드 분류
# ═══════════════════════════════════════════════════════════════

def _get_title(slide):
    if slide.shapes.title:
        t = slide.shapes.title.text.strip()
        if t:
            return t
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                return t[:300]
    return ''


OPENING_KEYWORDS = [
    'technical update briefing', 'azure tub', 'unlocking innovations',
    'updates by azure product', 'welcome', 'agenda', '목차',
]

CLOSEOUT_KEYWORDS = [
    'thank', 'q&a', 'questions', '감사합니다', '질문', 'wrap',
]


def classify_slide(title, selected, cat_areas):
    """
    Returns: (section, display_title)
      section: 'Opening' | 'IH' | 'CH' | 'YJ' | 'MH' | 'Closeout' | None(삭제)
    """
    if not title:
        return None, title

    tl_norm = _normalize(title)
    tl_lower = title.lower()

    # ── Opening (표지 ~ Agenda) ──
    for kw in OPENING_KEYWORDS:
        if kw in tl_lower:
            return 'Opening', title

    # ── Closeout (감사 슬라이드) ──
    for kw in CLOSEOUT_KEYWORDS:
        if kw in tl_lower:
            return 'Closeout', title

    # ── 선정 항목 exact match ──
    if tl_norm in selected:
        return selected[tl_norm]['area'], title

    # ── 카테고리 헤더 ──
    title_clean = ' '.join(title.split())   # trailing whitespace 정리
    if title_clean in KNOWN_CATEGORIES:
        if title_clean in cat_areas:         # 해당 카테고리에 선정 항목 있음
            return cat_areas[title_clean], title
        else:
            return None, title              # 선정 항목 없는 카테고리 → 삭제

    # ── 미매칭 → 삭제 ──
    return None, title


# ═══════════════════════════════════════════════════════════════
# 4. 슬라이드 삭제
# ═══════════════════════════════════════════════════════════════

def delete_slides(prs, indices):
    sldIdLst = prs.slides._sldIdLst
    all_ids  = list(sldIdLst)
    for idx in sorted(indices, reverse=True):
        el  = all_ids[idx]
        rId = el.get(f'{{{REL_NS}}}id')
        sldIdLst.remove(el)
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


# ═══════════════════════════════════════════════════════════════
# 5. 슬라이드 재정렬
# ═══════════════════════════════════════════════════════════════

def reorder_slides(prs, new_order):
    """new_order: 현재 인덱스 목록 (원하는 순서로)"""
    sldIdLst = prs.slides._sldIdLst
    saved    = list(sldIdLst)
    for el in saved:
        sldIdLst.remove(el)
    for idx in new_order:
        sldIdLst.append(saved[idx])


# ═══════════════════════════════════════════════════════════════
# 6. Section XML 삽입
# ═══════════════════════════════════════════════════════════════

def add_sections(prs, section_map):
    """
    section_map: [(section_name, [slide_index, ...]), ...]
    slide_index는 최종 정렬 후의 0-based 인덱스
    """
    prs_elem = prs.slides._sldIdLst.getparent()

    # 기존 sectionLst 제거
    for old in prs_elem.findall(f'.//{{{P14_NS}}}sectionLst'):
        parent = old.getparent()
        if parent is not None:
            parent.getparent().remove(parent)

    # extLst 확보
    extLst = prs_elem.find(qn('p:extLst'))
    if extLst is None:
        extLst = etree.SubElement(prs_elem, qn('p:extLst'))

    ext_el = etree.SubElement(extLst, qn('p:ext'))
    ext_el.set('uri', EXT_URI)

    sectionLst = etree.SubElement(
        ext_el, f'{{{P14_NS}}}sectionLst',
        nsmap={'p14': P14_NS},
    )

    # 슬라이드 ID 목록 (재정렬 완료 후)
    slide_ids = [prs.slides._sldIdLst[i].get('id') for i in range(len(prs.slides))]

    for name, indices in section_map:
        if not indices:
            continue
        sec = etree.SubElement(sectionLst, f'{{{P14_NS}}}section')
        sec.set('name', name)
        sec.set('id', '{' + str(uuid.uuid4()).upper() + '}')
        inner = etree.SubElement(sec, f'{{{P14_NS}}}sldIdLst')
        for idx in indices:
            sldId_el = etree.SubElement(inner, f'{{{P14_NS}}}sldId')
            sldId_el.set('id', slide_ids[idx])


# ═══════════════════════════════════════════════════════════════
# 7. 메인 처리
# ═══════════════════════════════════════════════════════════════

def process(pptx_path, excel_path, sheet_name=None, output_path=None):
    if output_path is None:
        base, ext = os.path.splitext(pptx_path)
        output_path = base + '_sectioned' + ext

    # ── Excel 로드 ──
    selected, cat_areas = load_voting_data(excel_path, sheet_name)
    print(f"[Excel] 선정항목 {len(selected)}개, 카테고리 {len(cat_areas)}개\n")

    # ── PPTX 로드 ──
    print(f"[PPTX] 로딩: {pptx_path}")
    prs   = Presentation(pptx_path)
    total = len(prs.slides)
    print(f"[PPTX] 전체 {total}장\n")

    # ── 분류 ──
    keep_map     = {}   # {original_index: section}
    delete_idx   = set()
    agenda_found = False

    print(f"{'슬라이드':>5}  {'섹션':^10}  제목")
    print('─' * 90)

    for i, slide in enumerate(prs.slides):
        title   = _get_title(slide)
        section, _ = classify_slide(title, selected, cat_areas)

        # Agenda 이후의 Opening 슬라이드는 그냥 삭제
        # (표지~Agenda 까지만 Opening으로 취급)
        if section == 'Opening' and 'agenda' in title.lower():
            agenda_found = True
        elif section == 'Opening' and agenda_found:
            # Agenda 이후에 또 Opening 키워드가 나오면 일반 삭제 처리
            section = None

        if section is None:
            delete_idx.add(i)
            print(f"  {i+1:>3}   {'[삭제]':^10}  {title[:70]}")
        else:
            keep_map[i] = section
            print(f"  {i+1:>3}   {('[' + section + ']'):^10}  {title[:70]}")

    print('─' * 90)
    print(f"유지 {len(keep_map)}장 / 삭제 {len(delete_idx)}장\n")

    # ── 삭제 ──
    delete_slides(prs, delete_idx)

    # ── 재정렬: Opening → IH → CH → YJ → MH → Closeout ──
    kept_sorted  = sorted(keep_map.keys())
    new_idx_map  = {orig: new for new, orig in enumerate(kept_sorted)}

    buckets = {s: [] for s in ['Opening'] + SECTION_ORDER + ['Closeout']}
    for orig in kept_sorted:
        sec   = keep_map[orig]
        new_i = new_idx_map[orig]
        bucket_key = sec if sec in buckets else 'MH'
        buckets[bucket_key].append(new_i)

    final_order = []
    for key in ['Opening'] + SECTION_ORDER + ['Closeout']:
        final_order.extend(buckets[key])

    reorder_slides(prs, final_order)

    # ── Section XML 삽입 (Opening·Closeout은 섹션 포함, 선택) ──
    pos = 0
    section_map = []
    for key in ['Opening'] + SECTION_ORDER + ['Closeout']:
        count = len(buckets[key])
        if count:
            section_map.append((key, list(range(pos, pos + count))))
            pos += count

    add_sections(prs, section_map)

    # ── 저장 ──
    prs.save(output_path)

    print('\n=== 섹션별 결과 ===')
    for key in ['Opening'] + SECTION_ORDER + ['Closeout']:
        cnt = len(buckets[key])
        if cnt:
            print(f"  {key:<12}: {cnt}장")
    print(f"\n저장 완료: {output_path}")
    print("PowerPoint 슬라이드 패널에서 Opening / IH / CH / YJ / MH / Closeout 섹션 확인")


# ═══════════════════════════════════════════════════════════════
# 8. CLI
# ═══════════════════════════════════════════════════════════════

if __name__ == '__main__':
    parser = argparse.ArgumentParser(
        description='TUB Voting Dashboard 기반 PPTX 섹션 자동 구성 에이전트'
    )
    parser.add_argument('pptx', help='처리할 PPTX 파일 경로')
    parser.add_argument(
        '--excel', default=None,
        help='Voting Dashboard Excel 경로 (기본: PPTX와 같은 폴더의 *Voting*Dashboard*.xlsx)'
    )
    parser.add_argument('--sheet', default=None, help='Excel 시트명 (기본: 최근 _result 시트 자동)')
    parser.add_argument('--output', default=None, help='출력 파일 경로 (기본: 원본명_sectioned.pptx)')
    args = parser.parse_args()

    if not os.path.exists(args.pptx):
        print(f"오류: PPTX 파일 없음 → {args.pptx}"); sys.exit(1)

    # Excel 자동 탐색: PPTX 폴더에서 Voting Dashboard 파일 찾기
    excel_path = args.excel
    if excel_path is None:
        folder = os.path.dirname(os.path.abspath(args.pptx))
        candidates = [
            f for f in os.listdir(folder)
            if f.endswith('.xlsx') and 'Voting' in f and 'Dashboard' in f
        ]
        if candidates:
            excel_path = os.path.join(folder, sorted(candidates)[-1])
            print(f"[Excel] 자동 탐색: {excel_path}")
        else:
            print("오류: Voting Dashboard Excel 파일을 찾을 수 없습니다.\n"
                  "       --excel 옵션으로 직접 경로를 지정하세요.")
            sys.exit(1)

    process(
        pptx_path   = args.pptx,
        excel_path  = excel_path,
        sheet_name  = args.sheet,
        output_path = args.output,
    )
